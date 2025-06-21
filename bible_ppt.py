import tkinter as tk
from tkinter import messagebox, PhotoImage
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import re
import sys
import os
import datetime

# 구절 파싱 함수
def parse_passage(passage):
    result = []
    for part in passage.split(","):
        part = part.strip()
        m = re.match(r"(\D+)\s*(\d+):(\d+)(?:-(\d+))?", part)
        if m:
            short_book, chapter, verse_start, verse_end = m.groups()
            book = book_alias.get(short_book, short_book)
            chapter = int(chapter)
            verse_start = int(verse_start)
            verse_end = int(verse_end) if verse_end else verse_start
            for verse in range(verse_start, verse_end + 1):
                result.append((book, chapter, verse))
    return result




def generate_ppt():
    input_text = entry.get()
    verses_to_find = parse_passage(input_text)

    prs = Presentation("src/template.pptx")
    missing = []
    layout = prs.slide_layouts[1]  

    for book, chapter, verse in verses_to_find:
        filtered = df[(df.iloc[:,0] == book) & 
                      (df.iloc[:,1] == chapter) & 
                      (df.iloc[:,2] == verse)]
        
        if filtered.empty:
            missing.append(f"{book} {chapter}:{verse}")
            continue

        for idx, row in filtered.iterrows():
            text = row[3]
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            content = slide.shapes.placeholders[10]

            title.text = ""
            content.text = f"{book}{chapter}:{verse} {text}"

    now = datetime.datetime.now()
    time_str = now.strftime("%m%d")
    file_extension = ".pptx"
    filename=f"{time_str}{file_extension}"
    prs.save(filename)
    
    if missing:
        msg = "일부 구절을 찾을 수 없습니다:\n" + "\n".join(missing)
        messagebox.showinfo("부분 완료", msg)
    else:
        messagebox.showinfo("완료", "PPT 생성이 완료되었습니다!")

def resource_path(relative_path):
    try:
        # PyInstaller로 묶였을 때
        base_path = sys._MEIPASS
    except Exception:
        # 일반 파이썬 실행 시
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)



# 성경 이름 축약어 매핑 (더 추가 가능)
df_alias = pd.read_excel("src/alias.xlsx", sheet_name="biblebookalias")
book_alias = dict(zip(df_alias.iloc[:, 0], df_alias.iloc[:, 1]))

df = pd.read_excel("src/resource.xlsx", sheet_name="bible")



# tkinter GUI 만들기
window = tk.Tk()
window.title("울산대성교회")

bible_img_path = resource_path("bible.png")
bible_img = PhotoImage(file=bible_img_path)

petra_img_path = resource_path("petra_one_third.png")
petra_img = PhotoImage(file=petra_img_path)

window.wm_iconphoto(True, bible_img)

label = tk.Label(window, text="구절 입력 (예: 창1:1-2, 민3:5) \n 성경이름 구절사이 띄어쓰기 주의 (창 1:1)",
                 image=petra_img,compound=tk.TOP)
label.pack(padx=10, pady=2)

entry = tk.Entry(window, width=30)
entry.pack(padx=10, pady=2)

button = tk.Button(window, text="G E N E R A T E", command=generate_ppt)
button.pack(padx=10, pady=2)

window.mainloop()